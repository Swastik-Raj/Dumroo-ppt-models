from __future__ import annotations

import math
from dataclasses import dataclass

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class FlowDiagramStyle:
    node_fill_rgb: tuple[int, int, int] = (234, 242, 255)
    node_line_rgb: tuple[int, int, int] = (79, 129, 189)
    text_rgb: tuple[int, int, int] = (31, 55, 94)
    font_name: str = "Calibri"
    font_size_pt: int = 16


def _set_rgb(color_format, rgb: tuple[int, int, int]) -> None:
    color_format.rgb = RGBColor(rgb[0], rgb[1], rgb[2])


def add_flow_diagram(slide, nodes: list[str], edges: list[tuple[str, str]], style: FlowDiagramStyle | None = None) -> None:
    style = style or FlowDiagramStyle()

    if not nodes:
        return

    # Standard slide dimensions
    slide_width = Inches(10.0)
    slide_height = Inches(7.5)

    # Calculate dynamic box sizes based on text length
    def calculate_box_size(text: str, num_nodes: int) -> tuple[float, float]:
        text_len = len(text)
        # Adjust max width based on number of nodes to fit more boxes
        if num_nodes > 6:
            max_width = 2.2
            base_width = 1.8
        else:
            max_width = 2.8
            base_width = 2.0

        width = max(base_width, min(max_width, base_width + (text_len - 10) * 0.03))
        height = 1.0 if text_len <= 20 else 1.1
        return Inches(width), Inches(height)

    # Adaptive layout based on number of nodes
    node_count = len(nodes)
    if node_count <= 4:
        max_per_row = node_count
    elif node_count <= 6:
        max_per_row = 3
    else:
        max_per_row = 4

    rows = (node_count + max_per_row - 1) // max_per_row

    # Calculate average box size for spacing
    avg_box_w = Inches(2.2)

    # Adaptive margins and gaps based on number of nodes
    if node_count > 6:
        left = Inches(0.5)
        top = Inches(2.0)
        hgap = Inches(0.4)
        vgap = Inches(0.6)
    else:
        left = Inches(0.8)
        top = Inches(2.2)
        hgap = Inches(0.5)
        vgap = Inches(0.7)

    node_to_shape = {}
    node_positions = {}

    # Calculate positions ensuring they fit within slide bounds
    for i, name in enumerate(nodes):
        r = i // max_per_row
        c = i % max_per_row
        box_w, box_h = calculate_box_size(name, node_count)

        # Calculate x position with spacing
        x = left + c * (avg_box_w + hgap)
        y = top + r * (Inches(1.1) + vgap)

        # Ensure box doesn't go off the right edge
        if x + box_w > slide_width - Inches(0.3):
            x = slide_width - box_w - Inches(0.3)

        # Ensure box doesn't go off the bottom edge
        if y + box_h > slide_height - Inches(0.3):
            y = slide_height - box_h - Inches(0.3)

        node_positions[name] = (x, y, box_w, box_h)

    # Draw arrows FIRST (so they appear behind boxes)
    for src, dst in edges:
        if src not in node_positions or dst not in node_positions:
            continue

        ax, ay, aw, ah = node_positions[src]
        bx, by, bw, bh = node_positions[dst]

        # Calculate center points
        acx = ax + aw // 2
        acy = ay + ah // 2
        bcx = bx + bw // 2
        bcy = by + bh // 2

        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, acx, acy, bcx, bcy)
        conn.line.width = Pt(2.5)
        _set_rgb(conn.line.color, style.node_line_rgb)

        # python-pptx 1.0.x doesn't expose arrowhead enums; simulate arrowheads
        # with a small rotated triangle shape placed near the destination.
        dx = float(bcx - acx)
        dy = float(bcy - acy)
        if abs(dx) + abs(dy) < 1.0:
            continue

        angle_deg = math.degrees(math.atan2(dy, dx))

        size = Inches(0.22)
        backoff = Inches(0.18)

        length = math.hypot(dx, dy)
        ux = dx / length
        uy = dy / length

        cx = bcx - int(ux * float(backoff))
        cy = bcy - int(uy * float(backoff))

        tri_left = int(cx - size / 2)
        tri_top = int(cy - size / 2)

        tri = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
            tri_left,
            tri_top,
            size,
            size,
        )
        tri.rotation = angle_deg + 90.0
        tri.fill.solid()
        _set_rgb(tri.fill.fore_color, style.node_line_rgb)
        tri.line.fill.background()

    # Now draw boxes on top of arrows
    for i, name in enumerate(nodes):
        x, y, box_w, box_h = node_positions[name]

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        fill = shape.fill
        fill.solid()
        _set_rgb(fill.fore_color, style.node_fill_rgb)

        shape.line.width = Pt(2.5)
        _set_rgb(shape.line.color, style.node_line_rgb)

        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = 1
        p = tf.paragraphs[0]
        p.alignment = 1
        run = p.add_run()

        # Truncate long node names
        node_text = name[:50] if len(name) > 50 else name
        run.text = node_text
        run.font.name = style.font_name

        # Auto-adjust font size for long text
        if len(name) > 25:
            run.font.size = Pt(14)
        elif len(name) > 15:
            run.font.size = Pt(15)
        else:
            run.font.size = Pt(style.font_size_pt)

        _set_rgb(run.font.color, style.text_rgb)

        node_to_shape[name] = shape
