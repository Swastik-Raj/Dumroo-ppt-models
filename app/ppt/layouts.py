from __future__ import annotations

from io import BytesIO

from PIL import Image, ImageEnhance
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.diagram.flow import FlowDiagramStyle, add_flow_diagram
from app.diagram.svg_flow import (SvgFlowStyle, render_flow_svg,
                                  svg_to_png_bytes)
from app.models import SlideSpec
from app.ppt.theme import Theme


def _set_rgb(color_format, rgb: tuple[int, int, int]) -> None:
    color_format.rgb = RGBColor(rgb[0], rgb[1], rgb[2])


def _is_dark(rgb: tuple[int, int, int]) -> bool:
    # Perceived luminance
    return (0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]) < 140


def _panel_colors(theme: Theme) -> tuple[tuple[int, int, int], tuple[int, int, int]]:
    if _is_dark(theme.background_rgb):
        return (17, 24, 39), (55, 65, 81)
    return (255, 255, 255), (229, 231, 235)


def _add_accent_bar(slide, theme: Theme) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(
            0), Inches(0), Inches(0.18), Inches(7.5)
    )
    bar.fill.solid()
    _set_rgb(bar.fill.fore_color, theme.accent_rgb)
    bar.line.fill.background()


def _add_title(slide, title: str, theme: Theme) -> None:
    tb = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.9))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()

    # Truncate very long titles
    title_text = title[:120] if len(title) > 120 else title
    run.text = title_text
    run.font.name = theme.font_title

    # Auto-reduce font size for long titles
    if len(title) > 60:
        run.font.size = Pt(32)
    elif len(title) > 40:
        run.font.size = Pt(36)
    else:
        run.font.size = Pt(40)

    _set_rgb(run.font.color, theme.title_rgb)
    try:
        p.space_after = Pt(4)
    except Exception:
        pass


def _add_body(slide, text: str, theme: Theme, left, top, width, height) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()

    # Truncate very long text to prevent overflow
    text_content = text[:800] if len(text) > 800 else text
    run.text = text_content
    run.font.name = theme.font_body

    # Auto-reduce font size for longer text
    if len(text) > 500:
        run.font.size = Pt(18)
    elif len(text) > 300:
        run.font.size = Pt(20)
    else:
        run.font.size = Pt(22)

    _set_rgb(run.font.color, theme.body_rgb)


def _add_bullets(slide, text: str, theme: Theme, left, top, width, height, *, font_size: int = 22) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)

    lines = [ln.strip() for ln in (text or "").split("\n") if ln.strip()]
    if len(lines) == 1:
        lines = [s.strip() for s in (text or "").replace(
            "•", "").split(".") if s.strip()]
    if not lines:
        lines = ["Key point", "Key point", "Key point"]

    # Auto-reduce font size if content is too long
    total_chars = sum(len(ln) for ln in lines)
    num_lines = len(lines)

    # More aggressive scaling for large content
    if total_chars > 1200 or num_lines > 10:
        font_size = 14
    elif total_chars > 1000 or num_lines > 8:
        font_size = 16
    elif total_chars > 800 or num_lines > 7:
        font_size = 18
    elif total_chars > 500 or num_lines > 6:
        font_size = 20

    # Dynamic bullet limits based on font size
    if font_size <= 14:
        max_bullets = 12
        max_chars = 150
        line_spacing = 4
    elif font_size <= 16:
        max_bullets = 10
        max_chars = 140
        line_spacing = 5
    elif font_size <= 18:
        max_bullets = 8
        max_chars = 130
        line_spacing = 6
    elif font_size <= 20:
        max_bullets = 7
        max_chars = 120
        line_spacing = 6
    else:
        max_bullets = 7
        max_chars = 110
        line_spacing = 6

    for i, line in enumerate(lines[:max_bullets]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        cleaned = line
        if cleaned.startswith("- "):
            cleaned = cleaned[2:].strip()
        if cleaned.startswith("• "):
            cleaned = cleaned[2:].strip()

        # Dynamic truncation based on font size
        if len(cleaned) > max_chars:
            cleaned = cleaned[:max_chars - 3] + "..."

        p.text = f"• {cleaned}"
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.font.name = theme.font_body
        p.font.size = Pt(font_size)
        _set_rgb(p.font.color, theme.body_rgb)
        try:
            p.space_after = Pt(line_spacing)
        except Exception:
            pass


def _center_crop_to_aspect(img: Image.Image, target_aspect: float) -> Image.Image:
    w, h = img.size
    aspect = w / h

    if abs(aspect - target_aspect) < 1e-3:
        return img

    if aspect > target_aspect:
        # too wide
        new_w = int(h * target_aspect)
        x0 = (w - new_w) // 2
        return img.crop((x0, 0, x0 + new_w, h))

    # too tall
    new_h = int(w / target_aspect)
    y0 = (h - new_h) // 2
    return img.crop((0, y0, w, y0 + new_h))


def _add_image_fit(slide, image_bytes: bytes, left, top, width, height) -> None:
    img = Image.open(BytesIO(image_bytes)).convert("RGB")
    target_aspect = float(width) / float(height)
    img = _center_crop_to_aspect(img, target_aspect)

    buf = BytesIO()
    img.save(buf, format="JPEG", quality=90)
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width=width, height=height)


def _darken_image_bytes(image_bytes: bytes, *, brightness: float = 0.55) -> bytes:
    # Make the hero image readable without relying on PPT transparency support.
    img = Image.open(BytesIO(image_bytes)).convert("RGB")
    enhancer = ImageEnhance.Brightness(img)
    img2 = enhancer.enhance(max(0.1, min(1.0, brightness)))
    buf = BytesIO()
    img2.save(buf, format="JPEG", quality=90)
    return buf.getvalue()


def _add_picture_bytes(slide, image_bytes: bytes, left, top, width, height) -> None:
    buf = BytesIO(image_bytes)
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width=width, height=height)


def layout_title_full_image(slide, slide_spec: SlideSpec, theme: Theme, image_bytes: bytes | None) -> None:
    _add_accent_bar(slide, theme)

    # Full-bleed image (if available) for a more "hero" look.
    if image_bytes:
        # Pre-darken the image so we don't rely on transparency support.
        hero_bytes = _darken_image_bytes(image_bytes, brightness=0.55)
        _add_image_fit(slide, hero_bytes, Inches(0.18),
                       Inches(0), Inches(13.15), Inches(7.5))

        # White title on image.
        tb = slide.shapes.add_textbox(
            Inches(0.9), Inches(1.8), Inches(11.4), Inches(1.4))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        run = p.add_run()

        # Truncate long titles
        title_text = slide_spec.title[:100] if len(slide_spec.title) > 100 else slide_spec.title
        run.text = title_text
        run.font.name = theme.font_title

        # Auto-adjust title size
        if len(slide_spec.title) > 60:
            run.font.size = Pt(42)
        elif len(slide_spec.title) > 40:
            run.font.size = Pt(48)
        else:
            run.font.size = Pt(54)

        _set_rgb(run.font.color, (255, 255, 255))

        subtitle = (slide_spec.content or "").strip()
        if subtitle:
            # Truncate long subtitles
            subtitle_text = subtitle[:250] if len(subtitle) > 250 else subtitle

            sb = slide.shapes.add_textbox(
                Inches(0.95), Inches(3.25), Inches(11.2), Inches(1.5))
            stf = sb.text_frame
            stf.clear()
            stf.word_wrap = True
            stf.vertical_anchor = MSO_ANCHOR.TOP
            sp = stf.paragraphs[0]
            srun = sp.add_run()
            srun.text = subtitle_text
            srun.font.name = theme.font_body

            # Auto-adjust subtitle size
            if len(subtitle) > 150:
                srun.font.size = Pt(20)
            else:
                srun.font.size = Pt(24)

            _set_rgb(srun.font.color, (229, 231, 235))
    else:
        _add_title(slide, slide_spec.title, theme)
        _add_body(slide, slide_spec.content, theme, Inches(
            0.9), Inches(1.5), Inches(11.6), Inches(5.5))


def layout_image_left_text_right(slide, slide_spec: SlideSpec, theme: Theme, image_bytes: bytes | None) -> None:
    _add_accent_bar(slide, theme)
    _add_title(slide, slide_spec.title, theme)

    panel_fill, panel_border = _panel_colors(theme)
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.55),
        Inches(1.55),
        Inches(6.2),
        Inches(5.6),
    )
    panel.fill.solid()
    _set_rgb(panel.fill.fore_color, panel_fill)
    panel.line.width = Pt(1)
    _set_rgb(panel.line.color, panel_border)

    if image_bytes:
        img_card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.75),
            Inches(1.55),
            Inches(5.6),
            Inches(5.6),
        )
        img_card.fill.solid()
        _set_rgb(img_card.fill.fore_color, panel_fill)
        img_card.line.width = Pt(1)
        _set_rgb(img_card.line.color, panel_border)

        _add_image_fit(slide, image_bytes, Inches(
            0.9), Inches(1.7), Inches(5.25), Inches(5.3))

    # Better text spacing / size with explicit bullets.
    _add_bullets(slide, slide_spec.content, theme, Inches(
        6.75), Inches(1.75), Inches(5.75), Inches(5.25), font_size=20)


def layout_diagram_center(slide, slide_spec: SlideSpec, theme: Theme, *, diagram_engine: str = "svg") -> None:
    _add_accent_bar(slide, theme)
    _add_title(slide, slide_spec.title, theme)
    if slide_spec.diagram:
        diagram_engine_l = (diagram_engine or "svg").strip().lower()

        diagram_png: bytes | None = None
        if diagram_engine_l == "svg":
            svg = render_flow_svg(
                slide_spec.diagram.nodes,
                slide_spec.diagram.edges,
                SvgFlowStyle(
                    background_rgb=theme.background_rgb,
                    node_fill_rgb=theme.diagram_node_fill_rgb,
                    node_line_rgb=theme.diagram_node_line_rgb,
                    text_rgb=theme.diagram_text_rgb,
                    font_family=theme.font_body,
                ),
                width_px=1600,
                height_px=700,
            )
            diagram_png = svg_to_png_bytes(svg, output_width_px=1600)

        panel_fill, panel_border = _panel_colors(theme)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.75),
            Inches(1.55),
            Inches(12.2),
            Inches(4.85),
        )
        card.fill.solid()
        _set_rgb(card.fill.fore_color, panel_fill)
        card.line.width = Pt(1)
        _set_rgb(card.line.color, panel_border)

        if diagram_png:
            _add_picture_bytes(slide, diagram_png, Inches(
                1.0), Inches(1.75), Inches(11.7), Inches(4.45))
        else:
            add_flow_diagram(
                slide,
                slide_spec.diagram.nodes,
                slide_spec.diagram.edges,
                FlowDiagramStyle(
                    node_fill_rgb=theme.diagram_node_fill_rgb,
                    node_line_rgb=theme.diagram_node_line_rgb,
                    text_rgb=theme.diagram_text_rgb,
                    font_name=theme.font_body,
                    font_size_pt=16,
                ),
            )
    # Small caption under diagram
    caption = (slide_spec.content or "").strip()
    if caption:
        # Truncate long captions
        caption_text = caption[:300] if len(caption) > 300 else caption
        _add_body(slide, caption_text, theme, Inches(0.9),
                  Inches(6.6), Inches(11.8), Inches(0.75))


def layout_bullets(slide, slide_spec: SlideSpec, theme: Theme) -> None:
    _add_accent_bar(slide, theme)
    _add_title(slide, slide_spec.title, theme)

    panel_fill, panel_border = _panel_colors(theme)
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(1.55),
        Inches(12.2),
        Inches(5.6),
    )
    panel.fill.solid()
    _set_rgb(panel.fill.fore_color, panel_fill)
    panel.line.width = Pt(1)
    _set_rgb(panel.line.color, panel_border)

    _add_bullets(slide, slide_spec.content, theme, Inches(
        1.0), Inches(1.7), Inches(11.7), Inches(5.35), font_size=22)
