from __future__ import annotations

import os
import re
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from app.models import PresentationSpec
from app.planner import SlidePlan
from app.ppt.layouts import (layout_bullets, layout_diagram_center,
                             layout_image_left_text_right,
                             layout_title_full_image)
from app.ppt.theme import Theme


def _safe_filename(name: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("_")
    return cleaned or "deck"


def _set_bg(slide, rgb: tuple[int, int, int]) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])


def build_pptx(
    presentation_spec: PresentationSpec,
    slide_plans: list[SlidePlan],
    image_search,
    output_dir: str,
    theme: Theme | None = None,
    diagram_engine: str = "svg",
) -> str:
    theme = theme or Theme()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    debug_images = (os.getenv("DEBUG_IMAGES", "").strip().lower()
                    in {"1", "true", "yes", "on"})

    for plan in slide_plans:
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, theme.background_rgb)

        image_bytes = None
        if plan.image_query and image_search is not None:
            try:
                if debug_images:
                    print(
                        f"[builder] image_query='{plan.image_query[:120]}' slide='{plan.slide.title}'")
                image_bytes = image_search.search_and_download(
                    plan.image_query)
                if debug_images:
                    if image_bytes:
                        print(
                            f"[builder] image_ok bytes={len(image_bytes)} slide='{plan.slide.title}'")
                    else:
                        print(
                            f"[builder] image_none slide='{plan.slide.title}'")
            except Exception:
                image_bytes = None
                if debug_images:
                    print(f"[builder] image_error slide='{plan.slide.title}'")
        elif debug_images and plan.image_query and image_search is None:
            print(
                f"[builder] image_search_disabled (no UNSPLASH_ACCESS_KEY) slide='{plan.slide.title}'")

        if plan.layout == "title_full_image":
            layout_title_full_image(slide, plan.slide, theme, image_bytes)
        elif plan.layout == "image_left_text_right":
            layout_image_left_text_right(slide, plan.slide, theme, image_bytes)
        elif plan.layout == "diagram_center":
            layout_diagram_center(slide, plan.slide, theme,
                                  diagram_engine=diagram_engine)
        else:
            layout_bullets(slide, plan.slide, theme)

    os.makedirs(output_dir, exist_ok=True)

    ts = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    filename = f"{_safe_filename(presentation_spec.title)}_{ts}.pptx"
    out_path = os.path.join(output_dir, filename)
    prs.save(out_path)
    return out_path
